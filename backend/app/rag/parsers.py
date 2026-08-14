"""文档解析器：txt/md/html/pdf/docx/xlsx/pptx → 带元数据的文本段。

每段返回 (text, page)；page 对 PDF 有效，其余为 None。
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import List, Tuple

SUPPORTED_TYPES = {
    "txt": "text", "md": "text", "markdown": "text", "html": "text", "htm": "text", "csv": "text",
    "pdf": "pdf", "docx": "docx", "xlsx": "xlsx", "pptx": "pptx",
}


class ParseError(RuntimeError):
    pass


def parse_file(path: Path) -> List[Tuple[str, int | None]]:
    """解析文档，返回 [(text, page_or_None), ...]。"""
    ext = path.suffix.lower().lstrip(".")
    if ext not in SUPPORTED_TYPES:
        raise ParseError(f"不支持的文件类型: .{ext}（支持: {', '.join(sorted(SUPPORTED_TYPES))}）")
    kind = SUPPORTED_TYPES[ext]
    try:
        if kind == "text":
            return _parse_text(path)
        if kind == "pdf":
            return _parse_pdf(path)
        if kind == "docx":
            return _parse_docx(path)
        if kind == "xlsx":
            return _parse_xlsx(path)
        if kind == "pptx":
            return _parse_pptx(path)
    except ParseError:
        raise
    except Exception as e:  # noqa: BLE001
        raise ParseError(f"解析失败 ({ext}): {e}") from e
    raise ParseError("未知文件类型")


def _clean(text: str) -> str:
    text = text.replace("\x00", "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _parse_text(path: Path) -> List[Tuple[str, int | None]]:
    raw = path.read_bytes()
    for enc in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return [(_clean(raw.decode(enc)), None)]
        except UnicodeDecodeError:
            continue
    return [(_clean(raw.decode("utf-8", errors="replace")), None)]


def _parse_pdf(path: Path) -> List[Tuple[str, int | None]]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: List[Tuple[str, int | None]] = []
    for i, page in enumerate(reader.pages, start=1):
        text = _clean(page.extract_text() or "")
        if text:
            parts.append((text, i))
    if not parts:
        raise ParseError("PDF 未提取到文本（可能是扫描件，请使用 OCR 版本）")
    return parts


def _parse_docx(path: Path) -> List[Tuple[str, int | None]]:
    from docx import Document

    doc = Document(str(path))
    lines: list[str] = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            lines.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
    if not lines:
        raise ParseError("docx 未提取到文本")
    return [(_clean("\n".join(lines)), None)]


def _parse_xlsx(path: Path) -> List[Tuple[str, int | None]]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        header = [f"{ws.title}"] if ws.title else []
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v).strip() for v in row]
            if any(vals):
                header.append(" | ".join(vals))
        if len(header) > 1:
            parts.append(_clean("\n".join(header)))
    if not parts:
        raise ParseError("xlsx 未提取到内容")
    return [(p, None) for p in parts]


def _parse_pptx(path: Path) -> List[Tuple[str, int | None]]:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        lines.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
        if lines:
            parts.append((_clean("\n".join(lines)), i))
    if not parts:
        raise ParseError("pptx 未提取到文本")
    return parts
